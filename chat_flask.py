import openai
from sentence_transformers import SentenceTransformer
import faiss
import json
import numpy as np
import requests
from bs4 import BeautifulSoup
from googlesearch import search
from flask import Flask, render_template, request, flash, send_file
import os
import tempfile

# Libraries to parse uploaded files
import PyPDF2
import docx
from pptx import Presentation

app = Flask(__name__)
app.secret_key = "your_secret_key"  # Change this to a secure secret in production

# -------------------------
# 1) OPENROUTER / OPENAI SETUP
# -------------------------
OPENROUTER_API_KEY = ""  # Replace with your actual API key

client = openai.OpenAI(
    base_url="https://openrouter.ai/api/v1",
    api_key=OPENROUTER_API_KEY,
)

# -------------------------
# 2) FAISS & EMBEDDINGS SETUP (Chatbot functionality)
# -------------------------
faiss_index = faiss.read_index("alibaba_faiss.index")
with open("alibaba_metadata.json", "r", encoding="utf-8") as f:
    metadata = json.load(f)

model = SentenceTransformer("all-MiniLM-L6-v2")

HEADERS = {
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                  "AppleWebKit/537.36 (KHTML, like Gecko) "
                  "Chrome/119.0.0.0 Safari/537.36"
}

def search_docs(query, top_k=6):
    query_embedding = model.encode([query], convert_to_numpy=True)
    faiss.normalize_L2(query_embedding)
    distances, indices = faiss_index.search(query_embedding, top_k)
    results = []
    for dist, idx in zip(distances[0], indices[0]):
        doc_id = str(idx)
        if doc_id in metadata and metadata[doc_id]["content"].strip() and dist < 4.0:
            results.append(metadata[doc_id])
    return results

def web_search_alibaba(query, num_results=3):
    """
    Searches Alibaba Cloud website for data.
    This version appends "Saudi Arabia" to the search query to capture region-specific info.
    """
    search_query = f"site:alibabacloud.com Saudi Arabia {query}"
    results = []
    try:
        for url in search(search_query, num_results=num_results, lang="en"):
            try:
                response = requests.get(url, headers=HEADERS, timeout=10)
                response.raise_for_status()
                soup = BeautifulSoup(response.text, "html.parser")
                paragraphs = [p.get_text(strip=True) for p in soup.find_all("p")]
                content = "\n".join(paragraphs)
                title_tag = soup.find("title")
                title = title_tag.text.strip() if title_tag else "No Title"
                if len(content) > 200:
                    results.append({"url": url, "title": title, "content": content})
            except Exception as e:
                print(f"❌ Error fetching {url}: {e}")
    except Exception as e:
        print(f"❌ Error during web search: {e}")
    return results

def parse_uploaded_file(file_storage):
    filename = file_storage.filename.lower()
    with tempfile.NamedTemporaryFile(delete=False) as tmp:
        tmp_name = tmp.name
        file_storage.save(tmp_name)
    text_content = ""
    try:
        if filename.endswith(".pdf"):
            with open(tmp_name, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text_content += page.extract_text() + "\n"
        elif filename.endswith(".docx"):
            doc = docx.Document(tmp_name)
            for para in doc.paragraphs:
                text_content += para.text + "\n"
        elif filename.endswith(".pptx"):
            prs = Presentation(tmp_name)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"
        elif filename.endswith(".txt"):
            with open(tmp_name, "r", encoding="utf-8", errors="ignore") as f:
                text_content = f.read()
    except Exception as e:
        print(f"❌ Error parsing file {filename}: {e}")
        text_content = ""
    finally:
        os.remove(tmp_name)
    return text_content

def generate_response(query, attachment_text=None):
    internal_docs = search_docs(query, top_k=5)
    web_docs = web_search_alibaba(query, num_results=3)
    combined_docs = internal_docs + web_docs
    if attachment_text:
        combined_docs.append({
            "title": "User Attachment",
            "content": attachment_text
        })
    if not combined_docs:
        return "**AI Response:**\n\nI couldn't find relevant Alibaba Cloud data."
    context = "\n\n".join([
        f"**{doc['title']}**:\n{doc['content'][:600]}"
        for doc in combined_docs
    ])
    prompt = f"""
    You are an expert in **Alibaba Cloud**, with in-depth knowledge of product features, pricing, and services—especially for the Saudi Arabian region.

    **Rules:**  
    - Answer the user's question **ONLY** using the provided documentation below.
    - Do not add extra information outside of what is provided.
    - Respond in **Markdown**.
    - Provide a **thorough, detailed** answer.

    **User Question:** {query}

    **Relevant Documentation:**
    {context}

    **Instructions:**  
    - Summarize key details with bullet points or headings.
    - Provide in-depth explanations and examples where relevant.
    - Ensure your answer is fact-based and detailed.
    - Focus on Saudi region product features, pricing, and services when applicable.
    """
    completion = client.chat.completions.create(
        extra_headers={"X-Title": "Alibaba Chatbot"},
        model="qwen/qwen2.5-vl-72b-instruct:free",
        messages=[{"role": "user", "content": prompt}],
        max_completion_tokens=2048,
        temperature=0.7
    )
    response_text = completion.choices[0].message.content.strip()
    return response_text

@app.route("/", methods=["GET", "POST"])
def index():
    response_text = ""
    if request.method == "POST":
        query = request.form.get("query", "")
        attachment = request.files.get("attachment")
        attachment_text = ""
        if attachment and attachment.filename:
            # For general files (PDF, DOCX, TXT); if PPTX is uploaded, flash an error.
            if attachment.filename.lower().endswith(".pptx"):
                flash("PPTX files are not supported in this chat. Please use a different file type.", "error")
            else:
                attachment_text = parse_uploaded_file(attachment)
                if not attachment_text:
                    flash("Could not parse file or unsupported format.", "error")
        response_text = generate_response(query, attachment_text)
    return render_template("index.html", response=response_text)

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=80, debug=True)